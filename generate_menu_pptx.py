"""14tsubo 限定アードベッグ メニュー PowerPoint生成スクリプト
HTMLからフルテキストを自動抽出してPPTXを生成する。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from html.parser import HTMLParser
from pathlib import Path
import re

# パス設定
SCRIPT_DIR = Path(__file__).parent
HTML_PATH = SCRIPT_DIR / "index.html"
BOTTLES_DIR = SCRIPT_DIR / "assets" / "bottles"
LOGO_PATH = SCRIPT_DIR / "assets" / "logo" / "1t.png"
OUTPUT_DIR = SCRIPT_DIR / "pdf"
OUTPUT_DIR.mkdir(exist_ok=True)

# カラー定義
BG_COLOR = RGBColor(0x05, 0x08, 0x05)
NEON = RGBColor(0x3D, 0xF5, 0x8A)
INK = RGBColor(0xEA, 0xF5, 0xEA)
INK_DIM = RGBColor(0x8F, 0xA7, 0x90)
GOLD_LINE = RGBColor(0x1A, 0x2E, 0x1A)


# ============ HTMLパーサー ============

class MenuExtractor(HTMLParser):
    """HTMLメニューからボトル情報を抽出するパーサー"""

    def __init__(self):
        super().__init__()
        self.bottles = []
        self.current_card = None
        self.in_card = False
        self.in_card_name = False
        self.in_meta = False
        self.in_desc = False
        self.in_price = False
        self.current_lang = None
        self.capture = False
        self.depth = 0
        self.card_img = None

    def handle_starttag(self, tag, attrs):
        d = dict(attrs)
        cls = d.get("class", "")

        if tag == "div" and "card" == cls.strip():
            self.in_card = True
            self.current_card = {"name": "", "meta": "", "price": "", "img": "", "ja": "", "en": "", "ko": ""}
        if self.in_card and tag == "div" and "card-img" in cls:
            pass
        if self.in_card and tag == "img" and "card-img" not in cls:
            src = d.get("src", "")
            if "bottles/" in src and not self.current_card.get("img"):
                self.current_card["img"] = src.split("/")[-1]
        if self.in_card and tag == "img":
            src = d.get("src", "")
            if "bottles/" in src:
                self.current_card["img"] = src.split("/")[-1]
        if tag == "div" and "card-name" in cls:
            self.in_card_name = True
        if tag == "div" and "meta" == cls.strip():
            self.in_meta = True
        if tag == "div" and "price" in cls:
            self.in_price = True
        if tag == "div" and "desc" == cls.strip():
            self.in_desc = True
        if self.in_desc and tag == "span" and "lang" in d:
            self.current_lang = d["lang"]
            self.capture = True

    def handle_endtag(self, tag):
        if tag == "span" and self.capture:
            self.capture = False
            self.current_lang = None
        if tag == "div" and self.in_card_name:
            self.in_card_name = False
        if tag == "div" and self.in_meta:
            self.in_meta = False
        if tag == "div" and self.in_price:
            self.in_price = False
        if tag == "div" and self.in_desc and self.current_card:
            self.in_desc = False
            self.bottles.append(self.current_card)
            self.current_card = {"name": "", "meta": "", "price": "", "img": "", "ja": "", "en": "", "ko": ""}
            self.in_card = False

    def handle_data(self, data):
        if not self.current_card:
            return
        text = data.strip()
        if not text:
            return
        if self.in_card_name and not self.in_desc:
            if self.current_card["name"]:
                self.current_card["name"] += " "
            self.current_card["name"] += text
        if self.in_meta:
            self.current_card["meta"] += text
        if self.in_price:
            self.current_card["price"] += text
        if self.capture and self.current_lang:
            self.current_card[self.current_lang] += text


def extract_bottles():
    """HTMLからボトル情報を抽出する"""
    with open(HTML_PATH, "r", encoding="utf-8") as f:
        html = f.read()
    parser = MenuExtractor()
    parser.feed(html)
    return parser.bottles


# ============ PPTX生成 ============

def add_slide_bg(slide):
    """スライド背景を黒に設定する"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=INK, bold=False, alignment=PP_ALIGN.LEFT):
    """テキストボックスを追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def create_title_slide(prs, lang):
    """タイトルスライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)

    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(1.5), Inches(1.0), Inches(7.0), Inches(2.5)
        )

    titles = {"ja": "限定アードベッグ コレクション", "en": "Limited Ardbeg Collection", "ko": "한정 아드벡 컬렉션"}
    subtitles = {"ja": "全17銘柄", "en": "17 Bottles", "ko": "17종"}

    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(1.0),
                 titles[lang], font_size=36, color=NEON, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.5),
                 subtitles[lang], font_size=16, color=INK_DIM, alignment=PP_ALIGN.CENTER)


def create_bottle_slide(prs, bottle, lang):
    """ボトル1本分のスライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)

    # ボトル画像
    img_path = BOTTLES_DIR / bottle["img"]
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(0.3), Inches(0.5), Inches(3.0), Inches(6.0)
        )

    # 銘柄名
    name = bottle["name"].replace("\n", " ")
    add_text_box(slide, Inches(3.6), Inches(0.4), Inches(6.0), Inches(0.8),
                 name, font_size=24, color=NEON, bold=True)

    # Meta
    add_text_box(slide, Inches(3.6), Inches(1.1), Inches(6.0), Inches(0.4),
                 bottle["meta"], font_size=10, color=INK_DIM)

    # 価格
    price = bottle["price"].replace("/30ml", "").strip()
    add_text_box(slide, Inches(3.6), Inches(1.5), Inches(6.0), Inches(0.5),
                 f"{price} / 30ml", font_size=22, color=NEON, bold=True)

    # テイスティングノート（フルテキスト）
    desc = bottle.get(lang, "")
    add_text_box(slide, Inches(3.6), Inches(2.3), Inches(6.0), Inches(4.5),
                 desc, font_size=12, color=INK_DIM)


def create_footer_slide(prs, lang):
    """フッタースライドを作成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)

    footer_lines = {
        "ja": [
            "ARdbEG & Fruit Cocktail / イッテンヨンツボ",
            "📍 恵比寿横丁",
            "🕐 金・土・祝前日 20:00–26:00 ／ 日・平日 20:00–25:00",
            "🍸 立ち飲み ／ チャージなし",
            "📸 @1.4tsubo",
            "※ 価格は1ショット（30ml）税込",
        ],
        "en": [
            "ARdbEG & Fruit Cocktail / Itten-Yon-Tsubo",
            "📍 Ebisu Yokocho, Tokyo",
            "🕐 Fri, Sat & Holidays Eve 8PM–2AM / Sun & Weekdays 8PM–1AM",
            "🍸 Standing bar / No cover charge",
            "📸 @1.4tsubo",
            "* Prices per single shot (30ml), tax included",
        ],
        "ko": [
            "ARdbEG & Fruit Cocktail / 잇텐욘쯔보",
            "📍 에비스 요코초, 도쿄",
            "🕐 금·토·공휴일 전날 오후8시–새벽2시 / 일·평일 오후8시–새벽1시",
            "🍸 스탠딩 바 / 커버 차지 없음",
            "📸 @1.4tsubo",
            "※ 가격은 1샷(30ml) 세금 포함",
        ],
    }

    y = Inches(2.0)
    for i, line in enumerate(footer_lines[lang]):
        sz = 20 if i == 0 else 14
        clr = INK if i == 0 else INK_DIM
        bld = i == 0
        add_text_box(slide, Inches(1.0), y, Inches(8.0), Inches(0.5),
                     line, font_size=sz, color=clr, bold=bld, alignment=PP_ALIGN.CENTER)
        y += Inches(0.6)


def generate_menu(lang, bottles):
    """指定言語のメニューPPTXを生成する"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, lang)

    for bottle in bottles:
        create_bottle_slide(prs, bottle, lang)

    create_footer_slide(prs, lang)

    output_path = OUTPUT_DIR / f"14tsubo_menu_{lang}.pptx"
    prs.save(str(output_path))
    print(f"✅ {output_path}")
    return output_path


if __name__ == "__main__":
    bottles = extract_bottles()
    print(f"📖 HTMLから{len(bottles)}本のボトル情報を抽出")

    for b in bottles:
        print(f"  - {b['name'][:30]}... | img={b['img']} | {b['price']}")

    print()
    for lang in ["ja", "en", "ko"]:
        generate_menu(lang, bottles)

    print(f"\n🎉 全3言語のメニュー生成完了！")
    print(f"📁 出力先: {OUTPUT_DIR}")
